import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import io

try:
    from lxml import etree

    LXML_AVAILABLE = True
except ImportError:
    LXML_AVAILABLE = False
    print(
        "Warning: lxml library not found. Table cell border styling cannot be copied."
    )

from pptx.shapes.autoshape import Shape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.table import Table


def copy_fill_properties(source, target):
    """Utility function to copy fill properties between shapes."""
    if not (hasattr(source, "fill") and hasattr(target, "fill")):
        return

    try:
        fill_type = source.fill.type

        # Set fill type first
        if fill_type == MSO_FILL_TYPE.BACKGROUND:
            target.fill.background()
        elif fill_type == MSO_FILL_TYPE.SOLID:
            target.fill.solid()

            # Copy RGB color if available
            if (
                hasattr(source.fill.fore_color, "rgb")
                and source.fill.fore_color.rgb is not None
            ):
                target.fill.fore_color.rgb = source.fill.fore_color.rgb
            # No theme color handling - simplified

        # Copy transparency
        if hasattr(source.fill, "transparency"):
            target.fill.transparency = source.fill.transparency
    except Exception as e:
        print(f"Error copying fill properties: {e}")


def copy_line_properties(source, target):
    """Utility function to copy line/border properties between shapes."""
    if not (hasattr(source, "line") and hasattr(target, "line")):
        return

    try:
        # Set line width
        if hasattr(source.line, "width"):
            target.line.width = source.line.width

        # Set dash style if available
        if hasattr(source.line, "dash_style"):
            target.line.dash_style = source.line.dash_style

        # Copy line fill properties
        if hasattr(source.line, "fill") and hasattr(target.line, "fill"):
            line_fill_type = source.line.fill.type

            if line_fill_type == MSO_FILL_TYPE.BACKGROUND:
                target.line.fill.background()
            elif line_fill_type == MSO_FILL_TYPE.SOLID:
                target.line.fill.solid()

                # Handle RGB color
                if (
                    hasattr(source.line.fill.fore_color, "rgb")
                    and source.line.fill.fore_color.rgb is not None
                ):
                    target.line.fill.fore_color.rgb = source.line.fill.fore_color.rgb
                # No theme color handling - simplified
    except Exception as e:
        print(f"Error copying line properties: {e}")


def copy_font_properties(source_font, target_font):
    """Utility function to copy font properties between runs."""
    if source_font is None or target_font is None:
        return

    # Copy basic font properties
    for attr in ["name", "size", "bold", "italic", "underline"]:
        if hasattr(source_font, attr):
            try:
                setattr(target_font, attr, getattr(source_font, attr))
            except Exception:
                pass

    # Copy color properties
    try:
        # Copy RGB color only
        if hasattr(source_font.color, "rgb") and source_font.color.rgb is not None:
            target_font.color.rgb = source_font.color.rgb
        # No theme color handling - simplified
    except Exception:
        pass


def create_powerpoint_report(
    summary_csv_path: str,
    template_pptx_path: str,
    output_pptx_path: str,
    template_slide_index: int = 1,
) -> str | None:
    """
    Generates a PowerPoint report from a summary CSV and a template PPTX.
    It creates one slide per influencer by duplicating and populating the template slide
    with data from the CSV. Returns the path of the generated file or None on error.
    """
    try:
        # --- 1. Load Data ---
        print(f"Loading summary data from: {summary_csv_path}")
        if not os.path.exists(summary_csv_path):
            print(f"Error: Summary CSV file not found at {summary_csv_path}")
            return None
        df_summary = pd.read_csv(summary_csv_path)
        print(f"Loaded CSV with columns: {df_summary.columns.tolist()}")

        # --- 2. Load Presentation Template ---
        print(f"Loading template presentation from: {template_pptx_path}")
        if not os.path.exists(template_pptx_path):
            print(f"Error: PowerPoint template file not found at {template_pptx_path}")
            return None
        prs = Presentation(template_pptx_path)
        print(f"Loaded presentation with {len(prs.slides)} slides")

        # --- 3. Identify Template Slide Layout ---
        if len(prs.slides) <= template_slide_index:
            print(
                f"Error: Template slide index ({template_slide_index}) is out of bounds..."
            )
            return None

        template_slide = prs.slides[template_slide_index]
        template_slide_layout = template_slide.slide_layout
        print(f"Using slide index {template_slide_index} as template.")

        # Store all template shapes for reference
        template_shapes = list(template_slide.shapes)

        # --- 4. Generate Slides for Each Influencer ---
        initial_slide_count = len(prs.slides)
        print(f"Generating slides for {len(df_summary)} influencers...")

        for index, row in df_summary.iterrows():
            influencer_handle = row.get("influencer_handle", "Unknown")
            print(f"  Processing influencer: {influencer_handle}")

            # Create a new slide with the same layout as the template
            new_slide = prs.slides.add_slide(template_slide_layout)

            # Process each shape from the template
            for template_shape in template_shapes:
                new_shape = None
                try:
                    # --- Process by shape type ---
                    if template_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        try:
                            # Try to get the placeholder by index
                            placeholder_idx = template_shape.placeholder_format.idx
                            new_shape = new_slide.placeholders[placeholder_idx]
                        except (AttributeError, KeyError, IndexError):
                            # If placeholder not found, create appropriate shape type
                            if template_shape.has_text_frame:
                                new_shape = new_slide.shapes.add_textbox(
                                    template_shape.left,
                                    template_shape.top,
                                    template_shape.width,
                                    template_shape.height,
                                )
                            elif template_shape.has_table:
                                # We'll handle tables separately
                                continue
                            else:
                                # Skip other placeholder types we can't handle
                                continue

                    elif template_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        new_shape = new_slide.shapes.add_shape(
                            template_shape.auto_shape_type,
                            template_shape.left,
                            template_shape.top,
                            template_shape.width,
                            template_shape.height,
                        )
                        # Copy adjustments (e.g., for rounded corners)
                        if (
                            hasattr(template_shape, "adjustments")
                            and len(template_shape.adjustments) > 0
                        ):
                            try:
                                # Iterate and assign adjustments individually
                                for i in range(len(template_shape.adjustments)):
                                    new_shape.adjustments[i] = (
                                        template_shape.adjustments[i]
                                    )
                            except Exception as adj_e:
                                print(
                                    f"      Warning: Could not copy adjustments: {adj_e}"
                                )

                    elif template_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        new_shape = new_slide.shapes.add_textbox(
                            template_shape.left,
                            template_shape.top,
                            template_shape.width,
                            template_shape.height,
                        )

                    elif template_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # Copy picture shapes (like logos)
                        image_stream = io.BytesIO(template_shape.image.blob)
                        new_shape = new_slide.shapes.add_picture(
                            image_stream,
                            template_shape.left,
                            template_shape.top,
                            width=template_shape.width,
                            height=template_shape.height,
                        )
                        # Skip further text/formatting for pictures
                        continue

                    else:
                        # Skip unsupported shapes
                        continue

                    if new_shape is None:
                        continue

                    # --- Copy appearance properties ---
                    # Copy fill, line and other formatting
                    copy_fill_properties(template_shape, new_shape)
                    copy_line_properties(template_shape, new_shape)

                    # Copy rotation
                    if hasattr(template_shape, "rotation"):
                        new_shape.rotation = template_shape.rotation

                    # --- Process text content ---
                    if template_shape.has_text_frame and new_shape.has_text_frame:
                        template_tf = template_shape.text_frame
                        target_tf = new_shape.text_frame
                        target_tf.clear()  # Clear any default text

                        # --- FIX START: Remove initial empty paragraph if it exists ---
                        if (
                            len(target_tf.paragraphs) > 0
                            and target_tf.paragraphs[0].text == ""
                        ):
                            try:
                                p_element = target_tf.paragraphs[
                                    0
                                ]._p  # Access internal element
                                target_tf._txBody.remove(
                                    p_element
                                )  # Remove the empty paragraph element
                            except Exception as e:
                                print(
                                    f"    Warning: Could not remove initial empty paragraph: {e}"
                                )
                        # --- FIX END ---

                        # Set text frame properties first
                        for prop in [
                            "margin_bottom",
                            "margin_left",
                            "margin_right",
                            "margin_top",
                            "vertical_anchor",
                            "word_wrap",
                            "auto_size",
                        ]:
                            if hasattr(template_tf, prop):
                                try:
                                    setattr(target_tf, prop, getattr(template_tf, prop))
                                except:
                                    pass

                        # Process each paragraph
                        for p_idx, template_p in enumerate(template_tf.paragraphs):
                            # Get the original text with placeholders
                            combined_text = "".join(r.text for r in template_p.runs)
                            modified_text = combined_text

                            # Special case for {{influencer}} which should map to influencer_handle
                            if "{{influencer}}" in modified_text:
                                modified_text = modified_text.replace(
                                    "{{influencer}}", influencer_handle
                                )

                            # Handle platform-specific placeholders
                            platform_prefixes = ["yt", "ig", "tt"]
                            metrics = [
                                "posts",
                                "impressions",
                                "reach",
                                "likes_comments",
                                "eng_rate",
                            ]

                            for prefix in platform_prefixes:
                                for metric in metrics:
                                    placeholder = f"{{{{{prefix}_{metric}}}}}"
                                    if placeholder in modified_text:
                                        column_name = f"{prefix}_{metric}"

                                        if column_name in row:
                                            value = row[column_name]

                                            # Format values according to metric type
                                            if metric == "eng_rate":
                                                # Format as percentage
                                                formatted_value = f"{value:.2%}"
                                            elif metric in [
                                                "impressions",
                                                "reach",
                                                "likes_comments",
                                            ]:
                                                # Format as numbers with commas
                                                formatted_value = f"{int(value):,}"
                                            else:
                                                # Default formatting
                                                formatted_value = str(value)

                                            modified_text = modified_text.replace(
                                                placeholder, formatted_value
                                            )

                            # Handle standard placeholders for any column in the data
                            for col in df_summary.columns:
                                tag = f"{{{{{col}}}}}"
                                if tag in modified_text:
                                    raw_value = row.get(col, "")
                                    if isinstance(raw_value, (int, float)):
                                        if col == "avg_engagement_rate":
                                            formatted_value = f"{raw_value:.2%}"
                                        elif (
                                            "impressions" in col
                                            or "reach" in col
                                            or "engagements" in col
                                            or "posts" in col
                                        ):
                                            formatted_value = f"{raw_value:,}"
                                        else:
                                            formatted_value = str(raw_value)
                                    else:
                                        formatted_value = str(raw_value)
                                    modified_text = modified_text.replace(
                                        tag, formatted_value
                                    )

                            # Create new paragraph and copy formatting
                            target_p = target_tf.add_paragraph()
                            target_p.alignment = template_p.alignment
                            target_p.level = template_p.level

                            # Copy paragraph spacing (skip space_before for the first paragraph)
                            if (
                                p_idx > 0
                                and hasattr(template_p, "space_before")
                                and template_p.space_before is not None
                            ):
                                target_p.space_before = template_p.space_before
                            if (
                                hasattr(template_p, "space_after")
                                and template_p.space_after is not None
                            ):
                                target_p.space_after = template_p.space_after
                            if (
                                hasattr(template_p, "line_spacing")
                                and template_p.line_spacing is not None
                            ):
                                target_p.line_spacing = template_p.line_spacing

                            # --- Handle text content and formatting ---
                            # Trim whitespace from the final text *before* adding runs
                            final_text = modified_text.strip()

                            if template_p.runs:
                                # Check if we should use a single run for the whole paragraph
                                if len(template_p.runs) == 1 or all(
                                    run.text == "" for run in template_p.runs[1:]
                                ):
                                    # Use a single run with the template formatting
                                    template_run = template_p.runs[0]
                                    new_run = target_p.add_run()
                                    new_run.text = final_text  # Assign trimmed text

                                    # Copy font formatting
                                    copy_font_properties(
                                        template_run.font, new_run.font
                                    )
                                else:
                                    # Multiple runs with different formatting
                                    # In this case we'll just use the text with the first run's formatting
                                    template_run = template_p.runs[0]
                                    new_run = target_p.add_run()
                                    new_run.text = final_text  # Assign trimmed text

                                    # Copy font formatting from the first run
                                    copy_font_properties(
                                        template_run.font, new_run.font
                                    )
                            else:
                                # No runs - just set text directly
                                target_p.text = final_text  # Assign trimmed text

                except Exception as shape_e:
                    print(f"    Error processing shape: {shape_e}")
                    continue  # Continue with the next shape

        # --- 5. Remove Template Slide ---
        # Get the rId of the template slide (slide at index template_slide_index)
        xml_slides = prs.slides._sldIdLst
        if len(xml_slides) > template_slide_index:
            try:
                slides = list(xml_slides)
                idx = template_slide_index  # template slide index
                rId = prs.slides._sldIdLst[idx].rId
                # Remove the slide
                prs.part.drop_rel(rId)
                xml_slides.remove(slides[idx])
                print(f"Removed template slide at index {template_slide_index}")
            except Exception as e:
                print(f"Error removing template slide: {e}")
                # It's not critical if this fails; the template slide will just remain in the ppt

        # --- 6. Save the Presentation ---
        # Create the output directory if it doesn't exist
        output_dir = os.path.dirname(output_pptx_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        prs.save(output_pptx_path)
        print(f"Presentation saved to {output_pptx_path}")
        print(
            f"Created {len(prs.slides) - initial_slide_count + 1} slides for {len(df_summary)} influencers"
        )

        return output_pptx_path

    except Exception as e:
        print(f"Error generating PowerPoint report: {e}")
        import traceback

        traceback.print_exc()
        return None
