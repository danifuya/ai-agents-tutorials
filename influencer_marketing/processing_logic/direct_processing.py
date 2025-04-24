"""
Direct processing of files in Google Drive without webhooks

This module provides functions to process campaign files directly from Google Drive
without relying on webhooks. It can be used for manual processing or testing.
"""

import os
import time
from google.auth import default
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, MediaFileUpload
import io
from dotenv import load_dotenv
import asyncio
import nest_asyncio
import pandas as pd

# Import the processing functions
from processing_logic.data_analyzer import create_campaign_summary
from processing_logic.report_generator import create_powerpoint_report

# Load environment variables
load_dotenv()
nest_asyncio.apply()

# Configuration constants
FOLDER_ID = os.getenv("GOOGLE_DRIVE_FOLDER_ID")
DOWNLOAD_PATH = "./data/"
CAMPAIGN_CONTENT_FILENAME = "campaign.csv"  # Updated to use the new filename
TEMPLATE_PPTX_PATH = "./reports/template/template.pptx"
OUTPUT_PPTX_PATH = "./reports/generated_report.pptx"


def get_drive_service():
    """Create and return a Google Drive service object using ADC."""
    try:
        # Use ADC - this will automatically use the credentials from gcloud auth application-default login
        credentials, project = default(scopes=["https://www.googleapis.com/auth/drive"])

        # Refresh token if needed
        if hasattr(credentials, "refresh"):
            credentials.refresh(Request())

        return build("drive", "v3", credentials=credentials)
    except Exception as e:
        print(f"Error creating Drive service: {e}")
        raise


async def download_file_from_drive(file_id, local_path):
    """Download a file from Google Drive by its ID and save it locally."""
    service = get_drive_service()
    request = service.files().get_media(fileId=file_id)

    os.makedirs(os.path.dirname(local_path), exist_ok=True)

    with open(local_path, "wb") as f:
        downloader = MediaIoBaseDownload(f, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            print(f"Download {int(status.progress() * 100)}%.")

    print(f"Downloaded file to {local_path}")
    return local_path


async def process_campaign_data():
    """Run the data analysis using the direct function."""
    try:
        # Only requires the campaign content CSV now
        campaign_path = os.path.join(DOWNLOAD_PATH, CAMPAIGN_CONTENT_FILENAME)
        if not os.path.exists(campaign_path):
            print(
                f"ERROR: Campaign content file not found at {campaign_path} for processing."
            )
            return False

        print("Running direct data analysis function...")
        # Call the direct function
        df_summary = create_campaign_summary(campaign_path)

        # Convert to DataFrame for saving to CSV (function already returns DataFrame)
        if df_summary is not None:
            # Define the output path in the reports folder
            output_path = "./reports/generated_campaign_summary.csv"
            # Ensure the reports directory exists
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            df_summary.to_csv(output_path, index=False)
            print(f"Summary saved locally to {output_path}")

            # Upload the summary CSV to Google Drive
            print(f"Uploading summary CSV to Google Drive...")
            summary_drive_id = await upload_to_drive(
                output_path, "generated_campaign_summary.csv"
            )
            if summary_drive_id:
                print(
                    f"Summary CSV uploaded successfully to Google Drive with ID: {summary_drive_id}"
                )
            else:
                print("Failed to upload summary CSV to Google Drive.")

            # Now run the reporting function to generate the PowerPoint
            await generate_report()

            return True
        else:
            print("Direct data analysis function failed or returned no summary.")
            return False

    except Exception as e:
        print(f"An error occurred while processing campaign data: {e}")
        import traceback

        traceback.print_exc()
        return False


async def generate_report():
    """Generate the PowerPoint report using the direct function."""
    try:
        # Get the summary CSV from the reports directory
        summary_csv_path = "./reports/generated_campaign_summary.csv"

        # Ensure the output directory exists
        os.makedirs(os.path.dirname(OUTPUT_PPTX_PATH), exist_ok=True)

        if not os.path.exists(summary_csv_path):
            print(f"Error: Summary CSV not found at {summary_csv_path}")
            return False

        if not os.path.exists(TEMPLATE_PPTX_PATH):
            print(f"Error: PowerPoint template not found at {TEMPLATE_PPTX_PATH}")
            # If template doesn't exist, try creating a simple one
            if not create_simple_template_if_missing(TEMPLATE_PPTX_PATH):
                return False  # Failed to create template

        print("Running direct PowerPoint generation function...")
        # Call the direct function
        generated_path = create_powerpoint_report(
            summary_csv_path=summary_csv_path,
            template_pptx_path=TEMPLATE_PPTX_PATH,
            output_pptx_path=OUTPUT_PPTX_PATH,
        )

        if generated_path:
            print(
                f"Report generation function finished successfully. Path: {generated_path}"
            )
            # Upload the generated report back to Google Drive
            print("Uploading report back to Google Drive...")
            report_file_id = await upload_to_drive(
                OUTPUT_PPTX_PATH, "generated_report.pptx"
            )

            if report_file_id:
                print(
                    f"Report uploaded successfully to Google Drive with ID: {report_file_id}"
                )
            else:
                print("Failed to upload report to Google Drive")
            return True
        else:
            print("Direct PowerPoint generation function failed.")
            return False

    except Exception as e:
        print(f"An error occurred while generating the report: {e}")
        import traceback

        traceback.print_exc()
        return False


# Function to upload a file back to Google Drive
async def upload_to_drive(file_path, file_name):
    """Upload a file to the watched Google Drive folder."""
    try:
        service = get_drive_service()

        file_metadata = {"name": file_name, "parents": [FOLDER_ID]}

        media = MediaFileUpload(file_path, resumable=True)

        file = (
            service.files()
            .create(body=file_metadata, media_body=media, fields="id")
            .execute()
        )

        print(
            f"Successfully uploaded {file_name} to Google Drive with ID: {file.get('id')}"
        )
        return file.get("id")

    except Exception as e:
        print(f"Error uploading file to Google Drive: {e}")
        return None


def create_simple_template_if_missing(template_path: str) -> bool:
    """Create a simple PowerPoint template if none exists."""
    try:
        from pptx import Presentation
        from pptx.util import Inches

        print(f"Template not found. Creating a simple template at {template_path}...")

        # Create directory if it doesn't exist
        os.makedirs(os.path.dirname(template_path), exist_ok=True)

        # Create a blank presentation
        prs = Presentation()

        # Add a title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Campaign Performance Report"
        subtitle.text = "Generated by Influencer Marketing Analysis Tool"

        # Add a content slide with placeholders
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Campaign Summary"

        # Save the presentation
        prs.save(template_path)
        print(f"Created simple template at {template_path}")
        return True

    except Exception as e:
        print(f"Error creating template: {e}")
        return False


async def list_files_in_folder():
    """List all files in the watched Google Drive folder."""
    try:
        service = get_drive_service()

        # Query for files in the specific folder
        results = (
            service.files()
            .list(
                q=f"'{FOLDER_ID}' in parents and trashed=false",
                fields="files(id, name, mimeType, createdTime)",
            )
            .execute()
        )

        items = results.get("files", [])

        if not items:
            print("No files found in the watched folder.")
            return []

        print(f"Found {len(items)} files in the watched folder:")
        for item in items:
            print(f"- {item['name']} (ID: {item['id']}, Type: {item['mimeType']})")

        return items

    except Exception as e:
        print(f"Error listing files: {e}")
        return []


async def main():
    """Main function to run the direct processing workflow."""
    try:
        # List files in the folder
        print("Checking for files in the watched folder...")
        files = await list_files_in_folder()

        if not files:
            print("No files found to process.")
            return

        # Check for the campaign file
        campaign_file = None
        for file in files:
            if file["name"] == CAMPAIGN_CONTENT_FILENAME:
                campaign_file = file
                break

        if not campaign_file:
            print(f"Campaign file '{CAMPAIGN_CONTENT_FILENAME}' not found.")
            return

        # Download the campaign file
        print(
            f"Found campaign file: {campaign_file['name']} (ID: {campaign_file['id']})"
        )
        campaign_path = os.path.join(DOWNLOAD_PATH, CAMPAIGN_CONTENT_FILENAME)
        await download_file_from_drive(campaign_file["id"], campaign_path)

        # Process the data
        print("Processing campaign data...")
        await process_campaign_data()

        print("Direct processing completed successfully.")

    except Exception as e:
        print(f"Error in main processing: {e}")
        import traceback

        traceback.print_exc()


# If this script is run directly
if __name__ == "__main__":
    # Create required directories
    os.makedirs(DOWNLOAD_PATH, exist_ok=True)
    os.makedirs(os.path.dirname(TEMPLATE_PPTX_PATH), exist_ok=True)
    os.makedirs(os.path.dirname(OUTPUT_PPTX_PATH), exist_ok=True)

    # Run the main async function
    asyncio.run(main())
